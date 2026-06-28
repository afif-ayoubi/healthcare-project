from __future__ import annotations

from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
PRES_DIR = ROOT / "presentation"
PRES_DIR.mkdir(parents=True, exist_ok=True)
CH = ROOT / "outputs" / "charts"

NAVY = "0B1F33"
TEAL = "0E7490"
MINT = "14B8A6"
CORAL = "E76F51"
GOLD = "F4A261"
SLATE = "64748B"
PALE = "EEF6F8"
LIGHT = "F7FAFC"
BORDER = "DDE7EC"
WHITE = "FFFFFF"
PURPLE = "7C3AED"
FONT = "Liberation Sans"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def rgb(hexstr: str) -> RGBColor:
    return RGBColor.from_string(hexstr)


def add_rect(slide, x, y, w, h, fill, line=None, radius=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, size=18, bold=False, color=NAVY, align=PP_ALIGN.LEFT,
             valign=MSO_ANCHOR.TOP, font=FONT, margin=0.03, italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear(); tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic; r.font.color.rgb = rgb(color)
    return box


def add_rich_text(slide, parts, x, y, w, h, size=18, color=NAVY, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame; tf.clear(); tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(.03); tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.alignment = align; p.space_after = Pt(0)
    for text, bold, clr in parts:
        r=p.add_run(); r.text=text; r.font.name=FONT; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=rgb(clr or color)
    return box


def add_bullets(slide, items, x, y, w, h, size=17, color=NAVY, bullet_color=TEAL, spacing=7):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame; tf.clear(); tf.word_wrap=True
    tf.margin_left=tf.margin_right=Inches(.04); tf.margin_top=tf.margin_bottom=Inches(.03)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text = item; p.level=0; p.font.name=FONT; p.font.size=Pt(size); p.font.color.rgb=rgb(color)
        p.space_after=Pt(spacing); p.bullet=True
    return box


def add_picture_fit(slide, path: Path, x, y, w, h, contain=True):
    im = Image.open(path)
    iw, ih = im.size
    box_ratio = w / h
    img_ratio = iw / ih
    if contain:
        if img_ratio > box_ratio:
            new_w = w; new_h = w / img_ratio; nx=x; ny=y+(h-new_h)/2
        else:
            new_h = h; new_w = h*img_ratio; nx=x+(w-new_w)/2; ny=y
        return slide.shapes.add_picture(str(path), nx, ny, width=new_w, height=new_h)
    pic = slide.shapes.add_picture(str(path), x, y, width=w, height=h)
    # crop to box ratio if needed
    if img_ratio > box_ratio:
        crop = (1 - box_ratio/img_ratio)/2
        pic.crop_left = pic.crop_right = crop
    elif img_ratio < box_ratio:
        crop = (1 - img_ratio/box_ratio)/2
        pic.crop_top = pic.crop_bottom = crop
    return pic


def add_header(slide, title, slide_no, kicker="PROJECT FINDING"):
    add_text(slide, kicker.upper(), Inches(.55), Inches(.22), Inches(4), Inches(.28), size=10, bold=True, color=TEAL)
    add_text(slide, title, Inches(.55), Inches(.52), Inches(11.8), Inches(.55), size=26, bold=True, color=NAVY)
    add_rect(slide, Inches(.55), Inches(1.12), Inches(12.2), Inches(.025), TEAL, radius=False)
    add_text(slide, f"{slide_no}/7", Inches(12.25), Inches(.22), Inches(.5), Inches(.25), size=10, bold=True, color=SLATE, align=PP_ALIGN.RIGHT)


def add_footer(slide, source="WHO-derived project snapshot; Global Tuberculosis Report 2025 context"):
    add_text(slide, source, Inches(.55), Inches(7.18), Inches(10.9), Inches(.2), size=8, color=SLATE)
    add_text(slide, "MSBA382 · Healthcare Analytics", Inches(11.2), Inches(7.18), Inches(1.55), Inches(.2), size=8, color=SLATE, align=PP_ALIGN.RIGHT)


def add_stat_card(slide, x, y, w, h, label, value, accent, white=False):
    fill = WHITE if white else LIGHT
    add_rect(slide, x, y, w, h, fill, BORDER)
    add_rect(slide, x, y, Inches(.07), h, accent, radius=False)
    add_text(slide, label, x+Inches(.17), y+Inches(.14), w-Inches(.25), Inches(.28), size=11, bold=True, color=SLATE)
    add_text(slide, value, x+Inches(.17), y+Inches(.48), w-Inches(.25), Inches(.42), size=25, bold=True, color=NAVY)


def set_notes(slide, text):
    tf = slide.notes_slide.notes_text_frame
    tf.text = text


prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

notes = []

# Slide 1
slide = prs.slides.add_slide(blank)
bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = rgb(NAVY)
add_text(slide, "MSBA382 · HEALTHCARE ANALYTICS", Inches(.72), Inches(.52), Inches(5), Inches(.35), size=12, bold=True, color="9EE7E5")
add_text(slide, "Global Tuberculosis\nBurden & Treatment Gaps", Inches(.72), Inches(1.06), Inches(8.5), Inches(1.5), size=34, bold=True, color=WHITE)
add_text(slide, "Trends, demographic disparities and geographic patterns", Inches(.75), Inches(2.55), Inches(8), Inches(.4), size=17, color="DCEFF4")
add_text(slide, "A decision-support dashboard for finding where burden and service-coverage gaps are concentrated.", Inches(.75), Inches(3.05), Inches(7.9), Inches(.72), size=20, color=WHITE)
add_stat_card(slide, Inches(.75), Inches(4.32), Inches(2.65), Inches(1.16), "ESTIMATED INCIDENT TB", "10.55M", CORAL, white=True)
add_stat_card(slide, Inches(3.58), Inches(4.32), Inches(2.65), Inches(1.16), "NOTIFIED CASES", "8.33M", TEAL, white=True)
add_stat_card(slide, Inches(6.41), Inches(4.32), Inches(2.65), Inches(1.16), "CALCULATED COVERAGE", "78.9%", MINT, white=True)
# right visual motif
for i, (xx, yy, rr, cc) in enumerate([
    (10.2,1.25,.52,CORAL),(11.2,1.85,.34,MINT),(10.55,2.65,.24,GOLD),(11.65,3.05,.45,TEAL),
    (10.0,3.75,.30,PURPLE),(11.15,4.15,.18,CORAL),(10.55,4.8,.42,MINT),(11.8,5.1,.28,GOLD)
]):
    shape=slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(xx), Inches(yy), Inches(rr), Inches(rr)); shape.fill.solid(); shape.fill.fore_color.rgb=rgb(cc); shape.line.fill.background()
add_text(slide, "Project snapshot · 2024", Inches(.75), Inches(5.76), Inches(3), Inches(.3), size=10, bold=True, color="9EE7E5")
add_text(slide, "WHO published rounded global headlines of 10.7M incident cases and 1.23M deaths; this dashboard aggregates the reproducible country-level extract.", Inches(.75), Inches(6.05), Inches(8.4), Inches(.55), size=10, color="C7D9E3")
add_text(slide, "[Student Name]", Inches(10.0), Inches(6.45), Inches(2.5), Inches(.3), size=12, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
add_text(slide, "23 June 2026", Inches(10.0), Inches(6.78), Inches(2.5), Inches(.25), size=10, color="C7D9E3", align=PP_ALIGN.RIGHT)
notes1 = (
"Tuberculosis is preventable and curable, yet it remains one of the world’s largest infectious-disease burdens. "
"My project turns WHO-derived country-level estimates into a Streamlit decision-support dashboard. In the reproducible 2024 snapshot, about 10.55 million people developed TB, 8.33 million new and relapse cases were notified, and calculated diagnosis-and-treatment coverage was 78.9%. The purpose is to show not only where TB is high, but where the health-system response may still be incomplete."
)
set_notes(slide, notes1); notes.append((1,"Opening",notes1))

# Slide 2
slide=prs.slides.add_slide(blank)
bg=slide.background.fill; bg.solid(); bg.fore_color.rgb=rgb(WHITE)
add_header(slide,"The analytical problem",2,"Why this dashboard")
# left cascade
add_text(slide,"A gap can emerge between disease burden and people reached by services",Inches(.65),Inches(1.35),Inches(6.1),Inches(.6),size=20,bold=True,color=NAVY)
add_rect(slide,Inches(.7),Inches(2.12),Inches(5.8),Inches(.9),PALE,BORDER)
add_text(slide,"Estimated incident TB",Inches(.95),Inches(2.28),Inches(2.8),Inches(.25),size=14,bold=True,color=SLATE)
add_text(slide,"10.55M",Inches(4.65),Inches(2.22),Inches(1.4),Inches(.45),size=28,bold=True,color=CORAL,align=PP_ALIGN.RIGHT)
add_rect(slide,Inches(.7),Inches(3.26),Inches(4.6),Inches(.9),"E8F4F6",BORDER)
add_text(slide,"Notified new & relapse",Inches(.95),Inches(3.42),Inches(2.8),Inches(.25),size=14,bold=True,color=SLATE)
add_text(slide,"8.33M",Inches(3.65),Inches(3.36),Inches(1.2),Inches(.45),size=28,bold=True,color=TEAL,align=PP_ALIGN.RIGHT)
# gap bracket/callout
add_rect(slide,Inches(.7),Inches(4.45),Inches(5.8),Inches(1.05),"FFF6ED","F2D7C7")
add_text(slide,"Estimated notification difference",Inches(.95),Inches(4.63),Inches(3.1),Inches(.28),size=14,bold=True,color=CORAL)
add_text(slide,"2.22M",Inches(4.55),Inches(4.58),Inches(1.45),Inches(.42),size=27,bold=True,color=NAVY,align=PP_ALIGN.RIGHT)
add_text(slide,"A surveillance signal — not proof every person was untreated",Inches(.95),Inches(5.02),Inches(4.9),Inches(.28),size=10,color=SLATE)
# right questions
add_text(slide,"Objective",Inches(7.15),Inches(1.38),Inches(2),Inches(.35),size=18,bold=True,color=TEAL)
add_text(slide,"Build one interactive tool that moves from description to prioritization.",Inches(7.15),Inches(1.78),Inches(5.4),Inches(.65),size=22,bold=True,color=NAVY)
q=["Where is TB burden highest?","How is it changing over time?","Who and which subtypes are most affected?","Where should health systems investigate and act first?"]
add_bullets(slide,q,Inches(7.25),Inches(2.75),Inches(5.1),Inches(2.6),size=17,spacing=11)
add_rect(slide,Inches(7.15),Inches(5.65),Inches(5.3),Inches(.75),PALE,BORDER)
add_text(slide,"Decision sequence: scale → geography → population → response",Inches(7.42),Inches(5.88),Inches(4.8),Inches(.28),size=14,bold=True,color=TEAL,align=PP_ALIGN.CENTER)
add_footer(slide)
notes2=(
"The analytical problem is the gap between estimated disease burden and the people who appear in routine services. "
"I calculate the notification difference as estimated incident cases minus notified new and relapse cases. It is deliberately labelled as an estimate and a surveillance signal, because it does not prove that every person in the difference was untreated. The dashboard answers four practical questions: where burden is highest, how trends are changing, who and which subtypes are most affected, and where a healthcare system should investigate first."
)
set_notes(slide,notes2); notes.append((2,"Problem and objective",notes2))

# Slide 3
slide=prs.slides.add_slide(blank); bg=slide.background.fill; bg.solid(); bg.fore_color.rgb=rgb(WHITE)
add_header(slide,"A reproducible, decision-ready data pipeline",3,"Data and method")
steps=[
    ("1","WHO-derived public snapshots","Burden, notifications, age/sex, TB–HIV, RR/MDR-TB",CORAL),
    ("2","Clean and integrate","Standardize ISO codes, country names, years and numeric fields",TEAL),
    ("3","Validate and derive","Unique country-year checks; rates, coverage, gaps and shares",MINT),
    ("4","Publish in Streamlit","Filters, maps, trends, demographics, downloads and forecast",GOLD),
]
xs=[.65,3.75,6.85,9.95]
for (num,title,desc,col),xx in zip(steps,xs):
    add_rect(slide,Inches(xx),Inches(1.52),Inches(2.72),Inches(2.18),LIGHT,BORDER)
    circ=slide.shapes.add_shape(MSO_SHAPE.OVAL,Inches(xx+.18),Inches(1.72),Inches(.52),Inches(.52)); circ.fill.solid(); circ.fill.fore_color.rgb=rgb(col); circ.line.fill.background()
    add_text(slide,num,Inches(xx+.18),Inches(1.80),Inches(.52),Inches(.28),size=16,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    add_text(slide,title,Inches(xx+.22),Inches(2.38),Inches(2.28),Inches(.55),size=17,bold=True,color=NAVY)
    add_text(slide,desc,Inches(xx+.22),Inches(3.00),Inches(2.25),Inches(.55),size=11,color=SLATE)
    if xx<9:
        add_text(slide,"→",Inches(xx+2.78),Inches(2.45),Inches(.32),Inches(.35),size=24,bold=True,color=TEAL,align=PP_ALIGN.CENTER)
# badges
badges=[("5.3K","country-year rows"),("217","countries / territories"),("25","years: 2000–2024"),("6","WHO regions")]
for i,(val,lab) in enumerate(badges):
    xx=.75+i*3.08
    add_rect(slide,Inches(xx),Inches(4.18),Inches(2.72),Inches(1.05),PALE,BORDER)
    add_text(slide,val,Inches(xx+.18),Inches(4.34),Inches(.92),Inches(.4),size=25,bold=True,color=TEAL)
    add_text(slide,lab,Inches(xx+1.02),Inches(4.42),Inches(1.48),Inches(.28),size=12,bold=True,color=NAVY)
# guardrail
add_rect(slide,Inches(.75),Inches(5.62),Inches(11.82),Inches(.85),"FFF6ED","F2D7C7")
add_text(slide,"Methodological guardrail",Inches(1.02),Inches(5.82),Inches(2.1),Inches(.28),size=14,bold=True,color=CORAL)
add_text(slide,"Country-level aggregates can differ slightly from WHO’s rounded report headlines; the latest WHO time series supersedes earlier releases.",Inches(3.05),Inches(5.80),Inches(9.1),Inches(.38),size=13,color=NAVY)
add_footer(slide,"WHO TB data; GTB-TME official repository; reproducible project manifest and hashes")
notes3=(
"The workflow is reproducible from raw input to dashboard. I retained the raw snapshots, standardized identifiers and numeric fields in pandas, merged the indicators into one country-year model, and recalculated global and regional rates from the country data. Automated checks confirm unique country-year records, the 2000-to-2024 range, non-negative burden values and male/female age-sex data. The final package contains the code, raw and processed files, a data dictionary, source hashes, and validation output."
)
set_notes(slide,notes3); notes.append((3,"Data and method",notes3))

# Slide 4
slide=prs.slides.add_slide(blank); bg=slide.background.fill; bg.solid(); bg.fore_color.rgb=rgb(WHITE)
add_header(slide,"Global progress is real — but the gap remains",4,"Core result")
add_rect(slide,Inches(.55),Inches(1.35),Inches(8.45),Inches(5.45),WHITE,BORDER)
add_picture_fit(slide,CH/"global_burden_notifications.png",Inches(.72),Inches(1.53),Inches(8.1),Inches(5.05),contain=True)
# right cards
cards=[
    ("2024 snapshot","10.55M estimated\n8.33M notified",CORAL),
    ("Remaining signal","2.22M notification\ndifference",TEAL),
    ("System shock","18.1% notification\nfall in 2020",GOLD),
]
for i,(lab,val,col) in enumerate(cards):
    yy=1.48+i*1.5
    add_rect(slide,Inches(9.25),Inches(yy),Inches(3.42),Inches(1.22),LIGHT,BORDER)
    add_rect(slide,Inches(9.25),Inches(yy),Inches(.08),Inches(1.22),col,radius=False)
    add_text(slide,lab.upper(),Inches(9.52),Inches(yy+.16),Inches(2.75),Inches(.23),size=10,bold=True,color=SLATE)
    add_text(slide,val,Inches(9.52),Inches(yy+.46),Inches(2.7),Inches(.58),size=19,bold=True,color=NAVY)
add_rect(slide,Inches(9.25),Inches(6.04),Inches(3.42),Inches(.72),PALE,BORDER)
add_text(slide,"The gap narrowed, but detection and linkage remain incomplete.",Inches(9.50),Inches(6.22),Inches(2.92),Inches(.32),size=12,bold=True,color=TEAL,align=PP_ALIGN.CENTER)
add_footer(slide)
notes4=(
"The core time series shows substantial improvement. Notifications more than doubled since 2000, and calculated coverage rose from 32.4% to 78.9%. However, the 2024 snapshot still contains a 2.22 million difference between estimated burden and notified cases. The 2020 disruption is especially important: notifications fell by 18.1% from 2019 while estimated incidence moved much less. This shows that service access and reporting can deteriorate rapidly even when the underlying disease burden does not disappear."
)
set_notes(slide,notes4); notes.append((4,"Global trend",notes4))

# Slide 5
slide=prs.slides.add_slide(blank); bg=slide.background.fill; bg.solid(); bg.fore_color.rgb=rgb(WHITE)
add_header(slide,"Priority depends on burden volume and service coverage",5,"Geography and gaps")
add_rect(slide,Inches(.55),Inches(1.35),Inches(8.5),Inches(5.52),WHITE,BORDER)
add_picture_fit(slide,CH/"country_notification_gaps_2024.png",Inches(.72),Inches(1.52),Inches(8.15),Inches(5.18),contain=True)
insights=[
    ("Largest absolute difference","Indonesia · 241K",CORAL),
    ("Largest total burden","India · 2.71M\n92% coverage",TEAL),
    ("Low coverage warning","Myanmar · 43.5%",GOLD),
    ("Regional perspective","732K regional gap",PURPLE),
]
for i,(lab,val,col) in enumerate(insights):
    yy=1.40+i*1.28
    add_rect(slide,Inches(9.28),Inches(yy),Inches(3.38),Inches(1.02),LIGHT,BORDER)
    circ=slide.shapes.add_shape(MSO_SHAPE.OVAL,Inches(9.48),Inches(yy+.22),Inches(.34),Inches(.34)); circ.fill.solid(); circ.fill.fore_color.rgb=rgb(col); circ.line.fill.background()
    add_text(slide,lab,Inches(9.95),Inches(yy+.14),Inches(2.38),Inches(.24),size=10,bold=True,color=SLATE)
    add_text(slide,val,Inches(9.95),Inches(yy+.44),Inches(2.38),Inches(.44),size=16,bold=True,color=NAVY)
add_rect(slide,Inches(9.28),Inches(6.32),Inches(3.38),Inches(.54),"EAF5F7",BORDER)
add_text(slide,"Use counts + rates + coverage together.",Inches(9.55),Inches(6.43),Inches(2.83),Inches(.28),size=12,bold=True,color=TEAL,align=PP_ALIGN.CENTER)
add_footer(slide)
notes5=(
"Geography changes the interpretation. Indonesia had the largest absolute notification difference, about 241 thousand, while India had the largest total burden but calculated coverage above 92%. Myanmar illustrates a different risk: a smaller absolute burden, but coverage below 44%. At regional level, the Western Pacific had the largest absolute notification difference. The dashboard therefore avoids ranking countries on one indicator alone and encourages users to combine absolute volume, population rate and service coverage."
)
set_notes(slide,notes5); notes.append((5,"Geography and gaps",notes5))

# Slide 6
slide=prs.slides.add_slide(blank); bg=slide.background.fill; bg.solid(); bg.fore_color.rgb=rgb(WHITE)
add_header(slide,"Different populations and clinical dimensions need different responses",6,"Demographics and subtypes")
add_rect(slide,Inches(.55),Inches(1.35),Inches(6.05),Inches(4.62),WHITE,BORDER)
add_rect(slide,Inches(6.75),Inches(1.35),Inches(6.03),Inches(4.62),WHITE,BORDER)
add_picture_fit(slide,CH/"age_sex_pyramid_2024.png",Inches(.72),Inches(1.52),Inches(5.7),Inches(4.25),contain=True)
add_picture_fit(slide,CH/"hiv_rr_burden_2024.png",Inches(6.92),Inches(1.52),Inches(5.68),Inches(4.25),contain=True)
# takeaway chips
chips=[
    ("Male : female ≈ 1.50","Tailor outreach to adult men",TEAL),
    ("South Africa · 134K TB–HIV","Integrate TB and HIV pathways",PURPLE),
    ("India · 130K RR/MDR-TB","Expand rapid resistance testing",CORAL),
]
for i,(head,desc,col) in enumerate(chips):
    xx=.55+i*4.08
    add_rect(slide,Inches(xx),Inches(6.20),Inches(3.85),Inches(.72),LIGHT,BORDER)
    add_rect(slide,Inches(xx),Inches(6.20),Inches(.07),Inches(.72),col,radius=False)
    add_text(slide,head,Inches(xx+.20),Inches(6.31),Inches(3.3),Inches(.20),size=11,bold=True,color=NAVY)
    add_text(slide,desc,Inches(xx+.20),Inches(6.54),Inches(3.3),Inches(.19),size=9.5,color=SLATE)
add_footer(slide)
notes6=(
"The age-sex analysis adds a population lens. Estimated male incidence was about one and a half times female incidence, with the largest adult burden in ages 25 to 54. The subtype analysis adds a clinical-programme lens. HIV-associated TB was concentrated in sub-Saharan Africa, especially South Africa, while drug-resistant TB followed a different pattern led by India and several European and Asian settings. These are not interchangeable problems: TB–HIV requires integrated testing and treatment pathways, while resistance requires rapid susceptibility testing and effective regimens."
)
set_notes(slide,notes6); notes.append((6,"Demographics and subtypes",notes6))

# Slide 7
slide=prs.slides.add_slide(blank); bg=slide.background.fill; bg.solid(); bg.fore_color.rgb=rgb(NAVY)
add_text(slide,"FROM INSIGHT TO ACTION",Inches(.7),Inches(.45),Inches(4),Inches(.28),size=11,bold=True,color="9EE7E5")
add_text(slide,"Four recommendations for a stronger TB response",Inches(.7),Inches(.82),Inches(9.5),Inches(.65),size=30,bold=True,color=WHITE)
recs=[
    ("1","Prioritize high-volume detection gaps","Focus case finding, diagnostics and reporting audits where absolute differences are largest.",CORAL),
    ("2","Protect continuity during disruption","Maintain diagnostics, drug supply, follow-up and data reporting during emergencies.",GOLD),
    ("3","Integrate TB–HIV and resistance pathways","Match the service model to co-epidemic and drug-resistance patterns.",PURPLE),
    ("4","Tailor outreach and govern the data","Use demographic targeting; investigate outliers and refresh with each WHO release.",MINT),
]
for i,(num,title,desc,col) in enumerate(recs):
    yy=1.78+i*1.03
    add_rect(slide,Inches(.72),Inches(yy),Inches(9.0),Inches(.82),"133149",None)
    circ=slide.shapes.add_shape(MSO_SHAPE.OVAL,Inches(.92),Inches(yy+.15),Inches(.48),Inches(.48)); circ.fill.solid(); circ.fill.fore_color.rgb=rgb(col); circ.line.fill.background()
    add_text(slide,num,Inches(.92),Inches(yy+.23),Inches(.48),Inches(.24),size=14,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    add_text(slide,title,Inches(1.58),Inches(yy+.12),Inches(3.6),Inches(.27),size=15,bold=True,color=WHITE)
    add_text(slide,desc,Inches(5.0),Inches(yy+.15),Inches(4.35),Inches(.44),size=11,color="D7E4EA")
# closing panel right
add_rect(slide,Inches(10.02),Inches(1.78),Inches(2.62),Inches(4.93),WHITE,None)
add_text(slide,"FIND",Inches(10.35),Inches(2.18),Inches(1.95),Inches(.38),size=25,bold=True,color=CORAL,align=PP_ALIGN.CENTER)
add_text(slide,"where burden and gaps are concentrated",Inches(10.35),Inches(2.62),Inches(1.95),Inches(.60),size=12,color=SLATE,align=PP_ALIGN.CENTER)
add_text(slide,"↓",Inches(10.75),Inches(3.25),Inches(1.15),Inches(.35),size=24,bold=True,color=TEAL,align=PP_ALIGN.CENTER)
add_text(slide,"TARGET",Inches(10.35),Inches(3.65),Inches(1.95),Inches(.38),size=25,bold=True,color=TEAL,align=PP_ALIGN.CENTER)
add_text(slide,"the population and programme response",Inches(10.35),Inches(4.10),Inches(1.95),Inches(.60),size=12,color=SLATE,align=PP_ALIGN.CENTER)
add_text(slide,"↓",Inches(10.75),Inches(4.73),Inches(1.15),Inches(.35),size=24,bold=True,color=MINT,align=PP_ALIGN.CENTER)
add_text(slide,"ACT",Inches(10.35),Inches(5.14),Inches(1.95),Inches(.38),size=25,bold=True,color=MINT,align=PP_ALIGN.CENTER)
add_text(slide,"with measurable, data-governed interventions",Inches(10.35),Inches(5.58),Inches(1.95),Inches(.60),size=12,color=SLATE,align=PP_ALIGN.CENTER)
add_text(slide,"Published dashboard URL: [PASTE STREAMLIT URL HERE]",Inches(.72),Inches(6.36),Inches(8.9),Inches(.25),size=11,bold=True,color="9EE7E5")
add_text(slide,"Thank you",Inches(.72),Inches(6.78),Inches(3),Inches(.28),size=17,bold=True,color=WHITE)
notes7=(
"The dashboard leads to four actions: prioritize high-volume detection gaps; protect TB services during disruption; integrate TB–HIV and drug-resistance pathways according to the local pattern; and tailor outreach while governing the data carefully. The value of the project is the sequence: find where burden and coverage gaps are concentrated, target the population and programme response, and act with measurable interventions. The app, data, code, manual, report and presentation are ready for deployment; the final step is publishing through the student’s Streamlit account."
)
set_notes(slide,notes7); notes.append((7,"Recommendations and close",notes7))

out=PRES_DIR/"TB_Healthcare_Analytics_Pitch.pptx"
prs.save(out)

# Save notes as a simple text file for reuse by the document generator.
lines=[]
for n,title,text in notes:
    lines += [f"SLIDE {n}: {title}", text, ""]
(PRES_DIR/"speaker_notes.txt").write_text("\n".join(lines),encoding="utf-8")
print(out)
